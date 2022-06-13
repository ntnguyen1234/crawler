from crawler.utils import *
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
try:
  import win32com.client
  import pythoncom
  from asposeslidescloud.apis.slides_api import SlidesApi
  from asposeslidescloud.models.export_format import ExportFormat
except Exception:
  pass

class Converter:
  def __init__(self, parameters: dict={}):
    self.parameters = parameters

  def initialize_ppt(self, current_folder, i: int, url: dict):
    return {
      'url'     : url['url'],
      'counter' : url['counter'],
      'title'   : url['title'],
      'date'    : url['date'],
      'location': current_folder.joinpath(f'{i+1} - {url["title"]}.pptx')
    }

  def process_ppt(self, urls_ppt, url: dict):
    ppt_info = url
    times_tried = 0
    while True:
      if times_tried == 2: return
      else:
        try:
          prs = Presentation(ppt_info['location'])
          break
        except (Exception, KeyError):
          print(f'\n\Presentation ==========================================\n\n{ppt_info["location"]}\n\n')
          print(traceback.format_exc())
          print('========================================================\n')
          time.sleep(5)
          times_tried += 1
    ppt_info['num_page'] = len(prs.slides)
    ppt_info['num_img']  = 0
    ppt_info['keywords'] = 0
    for slide in prs.slides:
      for shape in slide.shapes:
        if shape.has_text_frame:
          for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
              ppt_info['keywords'] += sum([run.text.lower().count(key.replace('"', '').lower()) for key in self.parameters['required']])
        types = [
          MSO_SHAPE_TYPE.CHART, 
          MSO_SHAPE_TYPE.DIAGRAM,
          MSO_SHAPE_TYPE.PICTURE,
        ]
        try:
          if shape.shape_type in types:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
              width, height = shape.image.size
              if width > 200 and height > 200: ppt_info['num_img'] += 1
            else: ppt_info['num_img'] += 1
        except Exception:
          print(f'\n\nget_ppt_shape ==========================================\n\n{ppt_info["url"]}\n\n')
          print(traceback.format_exc())
          print('========================================================\n')
    ppt_info['img_ratio'] = ppt_info['num_img']/ppt_info['num_page']
    urls_ppt['ppt'].append(ppt_info)

  def convert_aspose(self, location, filetype, content=None, out_file=None):
    slidesApi = SlidesApi(app_sid=self.parameters['aspose']['id_key'], app_key=self.parameters['aspose']['secret'])
    if content == None:
      with open(location, 'rb') as f:
        content = f.read()
    if filetype == 'pdf':
      response = slidesApi.convert(content, ExportFormat.PDF)
    elif filetype == 'pptx':
      response = slidesApi.convert(content, ExportFormat.PPTX)
    if out_file == None: out_file = location
    return Path(response).rename(out_file)
  
  def convert_winppt(self, location, ext, out_file=None, delay: int=0):
    file_type  = {
      'pdf' : 32,
      'pptx': 24,
    }
    if out_file == None:
      out_file = location.parent / location.stem
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    out = powerpoint.Presentations.Open(location, WithWindow=False)
    time.sleep(delay)
    out.SaveAs(out_file, file_type[ext])
    out.Close()
    powerpoint.Quit()
    return

  def convert_libre(self, folder, filetype):
    while True:
      if filetype == 'pptx':
        return subprocess.check_output(f'soffice --headless --convert-to {filetype} ./*.ppt', shell=True, cwd=folder)
      elif filetype == 'pdf':
        return subprocess.check_output(f'soffice --headless --convert-to {filetype} ./*.pptx', shell=True, cwd=folder)

  def convert(self, location, filetype, content=None, out_file=None):
    try:
      if filetype == 'pptx':
        ppt_path = location.parent / f'{location.stem}.ppt'
        with open(ppt_path, 'wb') as fb:
          fb.write(content)
      else: ppt_path = location
      return
      # return self.convert_winppt(ppt_path, filetype, out_file)
    except (Exception, pythoncom.com_error):
      print(f'\n\nwin32/libreoffice ==========================================\n\n{location}\n\n')
      print(traceback.format_exc())
      print('========================================================\n')
      return self.convert_aspose(location, filetype, content, out_file)

  def readwrite_ppt(self, content, content_text, current_folder, i: int, url: dict):
    ppt_info = self.initialize_ppt(current_folder, i, url)
    if content_text.startswith('PK'):
      with open(ppt_info['location'], 'wb') as fb:
        fb.write(content)
    else:
      self.convert(ppt_info['location'], 'pptx', content)
    return ppt_info
    # return self.process_ppt(ppt_info, url)